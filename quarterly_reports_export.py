"""
Quarterly Report Export - PDF and PowerPoint generators
"""
from io import BytesIO
from datetime import datetime
import tempfile
import os

# PDF Generation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_LEFT
from reportlab.graphics.shapes import Drawing, Rect, String, Line
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.piecharts import Pie

# PowerPoint Generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# Arabic text support
try:
    from arabic_reshaper import reshape
    from bidi.algorithm import get_display
    ARABIC_SUPPORT = True
except:
    ARABIC_SUPPORT = False

def format_arabic(text):
    """Format Arabic text for proper display"""
    if not text:
        return ""
    if ARABIC_SUPPORT:
        try:
            reshaped = reshape(str(text))
            return get_display(reshaped)
        except:
            return str(text)
    return str(text)

def format_currency(amount):
    """Format number as currency"""
    try:
        return f"{float(amount):,.2f}"
    except:
        return "0.00"

# ==================== PDF Generation ====================

def generate_quarterly_report_pdf(report_data: dict, comparison_data: dict = None, company_settings: dict = None):
    """Generate PDF for quarterly report"""
    buffer = BytesIO()
    
    # Create document
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=1.5*cm,
        leftMargin=1.5*cm,
        topMargin=2*cm,
        bottomMargin=2*cm
    )
    
    # Styles
    styles = getSampleStyleSheet()
    
    # Custom styles for Arabic
    title_style = ParagraphStyle(
        'ArabicTitle',
        parent=styles['Title'],
        fontSize=24,
        alignment=TA_CENTER,
        spaceAfter=20
    )
    
    heading_style = ParagraphStyle(
        'ArabicHeading',
        parent=styles['Heading2'],
        fontSize=16,
        alignment=TA_RIGHT,
        spaceAfter=12,
        textColor=colors.HexColor('#059669')
    )
    
    normal_style = ParagraphStyle(
        'ArabicNormal',
        parent=styles['Normal'],
        fontSize=11,
        alignment=TA_RIGHT,
        spaceAfter=6
    )
    
    elements = []
    
    # Get data
    period = report_data.get('period', {})
    summary = report_data.get('summary', {})
    totals = report_data.get('totals', {})
    chart_data = report_data.get('chart_data', [])
    top_revenue = report_data.get('top_revenue_accounts', [])
    top_expense = report_data.get('top_expense_accounts', [])
    
    quarter_names = {1: 'الأول', 2: 'الثاني', 3: 'الثالث', 4: 'الرابع'}
    quarter_name = quarter_names.get(period.get('quarter', 1), '')
    
    # Title
    title_text = format_arabic(f"التقرير الربعي - الربع {quarter_name} {period.get('year', '')}")
    elements.append(Paragraph(title_text, title_style))
    elements.append(Spacer(1, 20))
    
    # Company info if available
    if company_settings:
        company_name = company_settings.get('company_name_ar', '')
        if company_name:
            elements.append(Paragraph(format_arabic(company_name), heading_style))
            elements.append(Spacer(1, 10))
    
    # Report date
    report_date = datetime.now().strftime('%Y-%m-%d')
    elements.append(Paragraph(format_arabic(f"تاريخ التقرير: {report_date}"), normal_style))
    elements.append(Spacer(1, 20))
    
    # Summary Section
    elements.append(Paragraph(format_arabic("ملخص الأداء المالي"), heading_style))
    
    # Summary table
    summary_data = [
        [format_arabic('البيان'), format_arabic('المبلغ (ر.س)')],
        [format_arabic('الرصيد المرحل'), format_currency(summary.get('carried_forward', {}).get('net', 0))],
        [format_arabic('إجمالي الإيرادات'), format_currency(totals.get('total_revenue', 0))],
        [format_arabic('إجمالي المصروفات'), format_currency(totals.get('total_expense', 0))],
        [format_arabic('صافي الربح'), format_currency(totals.get('net_profit', 0))],
        [format_arabic('هامش الربح'), f"{totals.get('profit_margin', 0)}%"],
        [format_arabic('الرصيد البنكي المتاح'), format_currency(summary.get('bank_balance', 0))],
    ]
    
    summary_table = Table(summary_data, colWidths=[10*cm, 6*cm])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#059669')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('FONTSIZE', (0, 1), (-1, -1), 11),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
        ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e5e7eb')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f9fafb')]),
    ]))
    elements.append(summary_table)
    elements.append(Spacer(1, 30))
    
    # Monthly breakdown
    if chart_data:
        elements.append(Paragraph(format_arabic("التفصيل الشهري"), heading_style))
        
        monthly_data = [[format_arabic('الشهر'), format_arabic('الإيرادات'), format_arabic('المصروفات'), format_arabic('الصافي')]]
        for month in chart_data:
            monthly_data.append([
                format_arabic(month.get('month', '')),
                format_currency(month.get('revenue', 0)),
                format_currency(month.get('expense', 0)),
                format_currency(month.get('net', 0))
            ])
        
        monthly_table = Table(monthly_data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
        monthly_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3b82f6')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e5e7eb')),
        ]))
        elements.append(monthly_table)
        elements.append(Spacer(1, 30))
    
    # Top Revenue Accounts
    if top_revenue:
        elements.append(Paragraph(format_arabic("أعلى مصادر الإيرادات"), heading_style))
        
        revenue_data = [[format_arabic('#'), format_arabic('الحساب'), format_arabic('المبلغ (ر.س)')]]
        for i, acc in enumerate(top_revenue[:5], 1):
            revenue_data.append([
                str(i),
                format_arabic(acc.get('account_name', '')),
                format_currency(acc.get('amount', 0))
            ])
        
        revenue_table = Table(revenue_data, colWidths=[1.5*cm, 10*cm, 4.5*cm])
        revenue_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#10b981')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (0, -1), 'CENTER'),
            ('ALIGN', (1, 0), (1, -1), 'RIGHT'),
            ('ALIGN', (2, 0), (2, -1), 'CENTER'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e5e7eb')),
        ]))
        elements.append(revenue_table)
        elements.append(Spacer(1, 20))
    
    # Top Expense Accounts
    if top_expense:
        elements.append(Paragraph(format_arabic("أعلى بنود المصروفات"), heading_style))
        
        expense_data = [[format_arabic('#'), format_arabic('الحساب'), format_arabic('المبلغ (ر.س)')]]
        for i, acc in enumerate(top_expense[:5], 1):
            expense_data.append([
                str(i),
                format_arabic(acc.get('account_name', '')),
                format_currency(acc.get('amount', 0))
            ])
        
        expense_table = Table(expense_data, colWidths=[1.5*cm, 10*cm, 4.5*cm])
        expense_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#ef4444')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (0, -1), 'CENTER'),
            ('ALIGN', (1, 0), (1, -1), 'RIGHT'),
            ('ALIGN', (2, 0), (2, -1), 'CENTER'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e5e7eb')),
        ]))
        elements.append(expense_table)
        elements.append(Spacer(1, 20))
    
    # Comparison Section
    if comparison_data and comparison_data.get('comparison'):
        elements.append(PageBreak())
        elements.append(Paragraph(format_arabic("مقارنة الأرباع"), heading_style))
        
        comp_data = [[format_arabic('الربع'), format_arabic('الإيرادات'), format_arabic('المصروفات'), format_arabic('الصافي')]]
        for q in comparison_data['comparison']:
            label = f"Q{q['quarter']} {q['year']}"
            if q.get('is_current'):
                label += " (الحالي)"
            comp_data.append([
                format_arabic(label),
                format_currency(q.get('revenue', 0)),
                format_currency(q.get('expense', 0)),
                format_currency(q.get('net_profit', 0))
            ])
        
        comp_table = Table(comp_data, colWidths=[4*cm, 4*cm, 4*cm, 4*cm])
        comp_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#6366f1')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('TOPPADDING', (0, 0), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e5e7eb')),
        ]))
        elements.append(comp_table)
    
    # Footer
    elements.append(Spacer(1, 40))
    footer_text = format_arabic(f"تم إنشاء هذا التقرير بواسطة نظام المحاسبة السعودي - {report_date}")
    elements.append(Paragraph(footer_text, ParagraphStyle('Footer', fontSize=9, alignment=TA_CENTER, textColor=colors.gray)))
    
    # Build PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

# ==================== PowerPoint Generation ====================

def generate_quarterly_report_pptx(report_data: dict, comparison_data: dict = None, company_settings: dict = None):
    """Generate PowerPoint presentation for quarterly report"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Get data
    period = report_data.get('period', {})
    summary = report_data.get('summary', {})
    totals = report_data.get('totals', {})
    chart_data = report_data.get('chart_data', [])
    top_revenue = report_data.get('top_revenue_accounts', [])
    top_expense = report_data.get('top_expense_accounts', [])
    
    quarter_names = {1: 'الأول', 2: 'الثاني', 3: 'الثالث', 4: 'الرابع'}
    quarter_name = quarter_names.get(period.get('quarter', 1), '')
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(5, 150, 105)  # Emerald
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"التقرير الربعي"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"الربع {quarter_name} - {period.get('year', '')}"
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Company name
    if company_settings and company_settings.get('company_name_ar'):
        company_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
        company_frame = company_box.text_frame
        company_para = company_frame.paragraphs[0]
        company_para.text = company_settings.get('company_name_ar', '')
        company_para.font.size = Pt(20)
        company_para.font.color.rgb = RGBColor(255, 255, 255)
        company_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "ملخص الأداء المالي"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(5, 150, 105)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Summary boxes
    box_data = [
        ("إجمالي الإيرادات", format_currency(totals.get('total_revenue', 0)), RGBColor(16, 185, 129)),
        ("إجمالي المصروفات", format_currency(totals.get('total_expense', 0)), RGBColor(239, 68, 68)),
        ("صافي الربح", format_currency(totals.get('net_profit', 0)), RGBColor(59, 130, 246)),
        ("الرصيد البنكي", format_currency(summary.get('bank_balance', 0)), RGBColor(139, 92, 246)),
    ]
    
    box_width = Inches(2.8)
    box_height = Inches(1.8)
    start_x = Inches(0.8)
    gap = Inches(0.3)
    
    for i, (label, value, color) in enumerate(box_data):
        x = start_x + i * (box_width + gap)
        
        # Box background
        box = slide.shapes.add_shape(1, x, Inches(1.5), box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Label
        label_box = slide.shapes.add_textbox(x, Inches(1.7), box_width, Inches(0.5))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = RGBColor(255, 255, 255)
        label_para.alignment = PP_ALIGN.CENTER
        
        # Value
        value_box = slide.shapes.add_textbox(x, Inches(2.3), box_width, Inches(0.8))
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = f"{value} ر.س"
        value_para.font.size = Pt(22)
        value_para.font.bold = True
        value_para.font.color.rgb = RGBColor(255, 255, 255)
        value_para.alignment = PP_ALIGN.CENTER
    
    # Additional metrics
    metrics_data = [
        ("الرصيد المرحل", format_currency(summary.get('carried_forward', {}).get('net', 0))),
        ("هامش الربح", f"{totals.get('profit_margin', 0)}%"),
        ("الرصيد الختامي", format_currency(summary.get('closing_balance', {}).get('net', 0))),
    ]
    
    for i, (label, value) in enumerate(metrics_data):
        y = Inches(3.8) + i * Inches(0.8)
        
        label_box = slide.shapes.add_textbox(Inches(3), y, Inches(4), Inches(0.4))
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(16)
        label_para.alignment = PP_ALIGN.RIGHT
        
        value_box = slide.shapes.add_textbox(Inches(7), y, Inches(3), Inches(0.4))
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = f"{value} ر.س" if "%" not in value else value
        value_para.font.size = Pt(16)
        value_para.font.bold = True
        value_para.alignment = PP_ALIGN.LEFT
    
    # Slide 3: Monthly Chart
    if chart_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "الإيرادات والمصروفات الشهرية"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(5, 150, 105)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Chart
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = [m.get('month', '') for m in chart_data]
        chart_data_obj.add_series('الإيرادات', [m.get('revenue', 0) for m in chart_data])
        chart_data_obj.add_series('المصروفات', [m.get('expense', 0) for m in chart_data])
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, 
            Inches(1), Inches(1.5), 
            Inches(11.333), Inches(5.5),
            chart_data_obj
        ).chart
        
        # Style chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
    
    # Slide 4: Top Accounts
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "أعلى الحسابات"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(5, 150, 105)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Revenue accounts
    if top_revenue:
        rev_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
        rev_frame = rev_title.text_frame
        rev_para = rev_frame.paragraphs[0]
        rev_para.text = "أعلى مصادر الإيرادات"
        rev_para.font.size = Pt(18)
        rev_para.font.bold = True
        rev_para.font.color.rgb = RGBColor(16, 185, 129)
        
        for i, acc in enumerate(top_revenue[:5]):
            y = Inches(1.8) + i * Inches(0.5)
            acc_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(6), Inches(0.4))
            acc_frame = acc_box.text_frame
            acc_para = acc_frame.paragraphs[0]
            acc_para.text = f"{i+1}. {acc.get('account_name', '')} - {format_currency(acc.get('amount', 0))} ر.س"
            acc_para.font.size = Pt(12)
    
    # Expense accounts
    if top_expense:
        exp_title = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(6), Inches(0.5))
        exp_frame = exp_title.text_frame
        exp_para = exp_frame.paragraphs[0]
        exp_para.text = "أعلى بنود المصروفات"
        exp_para.font.size = Pt(18)
        exp_para.font.bold = True
        exp_para.font.color.rgb = RGBColor(239, 68, 68)
        
        for i, acc in enumerate(top_expense[:5]):
            y = Inches(1.8) + i * Inches(0.5)
            acc_box = slide.shapes.add_textbox(Inches(7), y, Inches(6), Inches(0.4))
            acc_frame = acc_box.text_frame
            acc_para = acc_frame.paragraphs[0]
            acc_para.text = f"{i+1}. {acc.get('account_name', '')} - {format_currency(acc.get('amount', 0))} ر.س"
            acc_para.font.size = Pt(12)
    
    # Slide 5: Comparison (if available)
    if comparison_data and comparison_data.get('comparison'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "مقارنة الأرباع"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(99, 102, 241)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Chart
        comp = comparison_data['comparison']
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = [f"Q{q['quarter']} {q['year']}" for q in comp]
        chart_data_obj.add_series('الإيرادات', [q.get('revenue', 0) for q in comp])
        chart_data_obj.add_series('المصروفات', [q.get('expense', 0) for q in comp])
        chart_data_obj.add_series('صافي الربح', [q.get('net_profit', 0) for q in comp])
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1.5),
            Inches(11.333), Inches(5.5),
            chart_data_obj
        ).chart
        
        chart.has_legend = True
    
    # Save to buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
